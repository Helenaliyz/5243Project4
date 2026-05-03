"""Generate the 10-minute Citi Bike project deck as Presentation.pptx."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).parent
IMG = ROOT / "Data Analysis"
OUT = ROOT / "Presentation.pptx"

# ---- Theme ----
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
TEAL = RGBColor(0x12, 0x86, 0x9E)
ACCENT = RGBColor(0xE8, 0x7A, 0x22)
LIGHT = RGBColor(0xF4, 0xF6, 0xF8)
DARK = RGBColor(0x1B, 0x1F, 0x24)
MUTED = RGBColor(0x55, 0x5F, 0x6D)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=LIGHT):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    return bg


def add_band(slide, color=NAVY, height=Inches(0.55)):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, height)
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = color
    return band


def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Calibri"
    run.font.color.rgb = color
    return tb


def add_bullets(slide, items, left, top, width, height,
                size=18, color=DARK, bullet_color=TEAL):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        bullet = p.add_run()
        bullet.text = "▸ "
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        bullet.font.name = "Calibri"
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def section_header(slide, kicker, title):
    add_band(slide, NAVY, Inches(0.55))
    add_text(slide, kicker, Inches(0.5), Inches(0.10), Inches(8), Inches(0.4),
             size=14, bold=True, color=LIGHT)
    add_text(slide, title, Inches(0.5), Inches(0.75), Inches(12.3), Inches(0.7),
             size=30, bold=True, color=NAVY)
    # accent rule
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(1.45),
                                  Inches(1.0), Inches(0.06))
    rule.line.fill.background()
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT


def add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, height=height)


def metric_card(slide, left, top, width, height, label, value, color=TEAL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 left, top, width, height)
    box.line.color.rgb = color
    box.line.width = Pt(1.25)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    add_text(slide, value, left, top + Inches(0.15), width, Inches(0.7),
             size=30, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, left, top + Inches(0.85), width, Inches(0.4),
             size=12, color=MUTED, align=PP_ALIGN.CENTER)


def footer(slide, n, total):
    add_text(slide, f"Citi Bike Demand  ·  STAT GR5243 Project 4",
             Inches(0.5), Inches(7.05), Inches(8), Inches(0.35),
             size=10, color=MUTED)
    add_text(slide, f"{n} / {total}",
             Inches(12.3), Inches(7.05), Inches(0.6), Inches(0.35),
             size=10, color=MUTED, align=PP_ALIGN.RIGHT)


TOTAL = 12

# =========================================================================
# Slide 1 — Title
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

# Decorative accent bars
bar1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          0, Inches(6.6), SW, Inches(0.12))
bar1.line.fill.background()
bar1.fill.solid(); bar1.fill.fore_color.rgb = ACCENT

bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          0, Inches(6.85), SW, Inches(0.05))
bar2.line.fill.background()
bar2.fill.solid(); bar2.fill.fore_color.rgb = TEAL

add_text(s, "STAT GR5243  ·  Project 4  ·  May 2026",
         Inches(0.7), Inches(1.0), Inches(12), Inches(0.5),
         size=16, bold=True, color=ACCENT)

add_text(s, "Predicting Daily Citi Bike Demand",
         Inches(0.7), Inches(1.7), Inches(12), Inches(1.4),
         size=48, bold=True, color=LIGHT)

add_text(s,
         "An end-to-end machine learning pipeline for\n"
         "station-level demand forecasting in New York City",
         Inches(0.7), Inches(3.2), Inches(12), Inches(1.6),
         size=22, color=LIGHT)

add_text(s, "Helena Li (yl6029)   ·   Wentao Zhong (wz2753)   ·   "
            "Kevin Ma (km4189)   ·   Ketaki Dabade (kvd2112)",
         Inches(0.7), Inches(5.6), Inches(12), Inches(0.5),
         size=16, color=LIGHT)

footer(s, 1, TOTAL)

# =========================================================================
# Slide 2 — The Problem
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "01  ·  MOTIVATION", "Why Predict Citi Bike Demand?")

add_bullets(s, [
    "Citi Bike is NYC's largest bike-share system — millions of trips per year.",
    "Operations team must rebalance bikes across stations every day.",
    "Bad forecasts → empty docks at commuter hubs, full docks near parks.",
    "Demand is driven by weather, calendar, geography AND station identity.",
    "Goal: predict the number of trips per station per day, one day ahead.",
], Inches(0.7), Inches(1.9), Inches(7.5), Inches(4.8), size=20)

# Right-side stat panel
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(8.7), Inches(1.9), Inches(4.1), Inches(4.5))
panel.line.color.rgb = TEAL
panel.line.width = Pt(1.5)
panel.fill.solid(); panel.fill.fore_color.rgb = LIGHT

add_text(s, "THE PREDICTIVE QUESTION",
         Inches(8.85), Inches(2.05), Inches(3.8), Inches(0.4),
         size=12, bold=True, color=TEAL)
add_text(s,
         "Given a station and a date,\nhow many trips will start there?",
         Inches(8.85), Inches(2.5), Inches(3.8), Inches(1.2),
         size=18, bold=True, color=NAVY)

add_text(s, "WHY IT'S HARD",
         Inches(8.85), Inches(3.9), Inches(3.8), Inches(0.4),
         size=12, bold=True, color=TEAL)
add_text(s,
         "• Right-skewed demand (long tail)\n"
         "• Weather × calendar × geography\n"
         "• Hundreds of heterogeneous stations\n"
         "• Strong temporal autocorrelation",
         Inches(8.85), Inches(4.35), Inches(3.8), Inches(2.0),
         size=14, color=DARK)

footer(s, 2, TOTAL)

# =========================================================================
# Slide 3 — Pipeline Overview
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "02  ·  APPROACH", "End-to-End Pipeline")

steps = [
    ("DATA", "Trips + Weather\n+ Station Context", NAVY),
    ("CLEAN", "Parse, dedupe,\nstandardize", TEAL),
    ("EDA", "Trends, distributions,\nweather effects", TEAL),
    ("CLUSTER", "K-Means on station\nprofiles (k=4)", ACCENT),
    ("FEATURES", "Lag, weather flags,\ncyclical calendar", ACCENT),
    ("MODEL", "Ridge / RF / GB\n+ tuned RF", NAVY),
]

box_w = Inches(1.95)
box_h = Inches(2.0)
gap = Inches(0.12)
total_w = box_w * len(steps) + gap * (len(steps) - 1)
start_left = (SW - total_w) // 2
top = Inches(2.4)

for i, (kicker, body, c) in enumerate(steps):
    left = start_left + i * (box_w + gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_w, box_h)
    box.line.color.rgb = c
    box.line.width = Pt(1.5)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    add_text(s, kicker, left, top + Inches(0.15), box_w, Inches(0.4),
             size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, body, left, top + Inches(0.7), box_w, Inches(1.2),
             size=14, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        arrow_left = left + box_w + gap // 2 - Inches(0.05)
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   arrow_left - Inches(0.04),
                                   top + box_h // 2 - Inches(0.08),
                                   Inches(0.18), Inches(0.18))
        arrow.line.fill.background()
        arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED

add_text(s,
         "Time-based 80/20 split  ·  metrics on the original (un-logged) trip-count scale",
         Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.5),
         size=16, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s,
         "Each step feeds the next — clustering becomes a feature, "
         "lag demand becomes the strongest predictor.",
         Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.5),
         size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

footer(s, 3, TOTAL)

# =========================================================================
# Slide 4 — Data Sources & Cleaning
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "03  ·  DATA", "Three Sources, One Analytic Table")

add_text(s, "SOURCES", Inches(0.7), Inches(1.95), Inches(6), Inches(0.4),
         size=14, bold=True, color=TEAL)
add_bullets(s, [
    "Citi Bike trips — monthly CSVs (Jan–Feb 2026).",
    "Daily NYC weather — Open-Meteo (temp, precip, wind).",
    "Station context — IDs, lat/lon, neighborhood flags (park, pier, plaza, square).",
], Inches(0.7), Inches(2.4), Inches(6), Inches(2.5), size=16)

add_text(s, "CLEANING", Inches(0.7), Inches(4.6), Inches(6), Inches(0.4),
         size=14, bold=True, color=TEAL)
add_bullets(s, [
    "Parse timestamps; drop trips with invalid station IDs.",
    "Standardize station names; verify coordinates.",
    "Group by (station, date) → daily_trip_count.",
    "Merge weather (on date) and station context (on station_id).",
], Inches(0.7), Inches(5.05), Inches(6), Inches(2.0), size=16)

# Right column: schema card
panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                           Inches(7.4), Inches(1.95), Inches(5.4), Inches(5.0))
panel.line.color.rgb = NAVY
panel.line.width = Pt(1.5)
panel.fill.solid(); panel.fill.fore_color.rgb = LIGHT

add_text(s, "ONE ROW PER STATION-DAY",
         Inches(7.6), Inches(2.1), Inches(5.0), Inches(0.4),
         size=14, bold=True, color=NAVY)
add_text(s,
         "station_id   ·   date   ·   daily_trip_count\n"
         "casual_share / member_share\n"
         "classic_bike_share / electric_bike_share\n"
         "mean_temperature   ·   precipitation_sum\n"
         "max_wind_speed\n"
         "start_lat   ·   start_lng\n"
         "has_park   ·   has_pier   ·   has_plaza   ·   has_square",
         Inches(7.6), Inches(2.6), Inches(5.0), Inches(4.0),
         size=14, color=DARK)

footer(s, 4, TOTAL)

# =========================================================================
# Slide 5 — EDA: Demand Patterns
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "04  ·  EDA", "Demand Has Strong Time & Weather Patterns")

add_image(s, IMG / "daily_demand_trend.png",
          Inches(0.5), Inches(1.85), width=Inches(7.0))
add_image(s, IMG / "weekday_weekend_average_total_trips.png",
          Inches(0.5), Inches(4.85), width=Inches(3.4))
add_image(s, IMG / "demand_distribution_histogram.png",
          Inches(4.1), Inches(4.85), width=Inches(3.4))

# Right column: takeaways
add_text(s, "WHAT WE SEE", Inches(7.9), Inches(1.85), Inches(5), Inches(0.4),
         size=14, bold=True, color=TEAL)
add_bullets(s, [
    "System ramps from <2k to ~48k trips/day in mid-January.",
    "Weekday avg ≈ 17.5k trips  ·  Weekend ≈ 10.8k (~60% premium).",
    "Station-day demand is heavily right-skewed → log-transform target.",
    "Median ≈ 40 trips/day, long tail past 400.",
], Inches(7.9), Inches(2.3), Inches(5.1), Inches(4.5), size=16)

footer(s, 5, TOTAL)

# =========================================================================
# Slide 6 — EDA: Weather & Geography
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "04  ·  EDA", "Weather Suppresses Demand; Geography Concentrates It")

add_image(s, IMG / "demand_vs_temperature.png",
          Inches(0.5), Inches(1.9), width=Inches(4.2))
add_image(s, IMG / "demand_vs_precipitation.png",
          Inches(0.5), Inches(4.55), width=Inches(4.2))
add_image(s, IMG / "top_10_busiest_stations.png",
          Inches(5.0), Inches(1.9), width=Inches(5.0))

add_text(s, "TAKEAWAYS", Inches(10.3), Inches(1.9), Inches(2.7), Inches(0.4),
         size=14, bold=True, color=TEAL)
add_bullets(s, [
    "Highest-demand days cluster at warmer temps.",
    "Days with >10 mm precipitation almost always low-demand.",
    "Top stations: Chelsea, Midtown, Union Sq, Pier 61.",
    "Mix of commuter and leisure hubs → motivates clustering.",
], Inches(10.3), Inches(2.35), Inches(2.7), Inches(4.5), size=13)

footer(s, 6, TOTAL)

# =========================================================================
# Slide 7 — Clustering
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "05  ·  UNSUPERVISED", "Stations Fall Into 4 Behavioral Clusters")

add_bullets(s, [
    "Aggregate to station-level profile (NOT station-day):",
    "    avg / median / max demand, demand volatility (CV),",
    "    casual vs. member share, electric-bike share,",
    "    weekend-to-weekday ratio, lat/lon, neighborhood flags.",
    "Standardize → K-Means.",
    "Choose k = 4 via elbow + silhouette.",
    "Visualize in 2D with PCA.",
], Inches(0.7), Inches(1.95), Inches(7), Inches(4.5), size=18)

# Cluster cards
cards = [
    ("Cluster A", "Commuter hubs", "high weekday volume,\nmember-dominated", NAVY),
    ("Cluster B", "Leisure / waterfront", "elevated casual share,\nweekend-heavy", ACCENT),
    ("Cluster C", "Mixed-use", "moderate demand\nacross rider types", TEAL),
    ("Cluster D", "Low-demand", "low average and\npeak demand", MUTED),
]
card_w, card_h = Inches(2.45), Inches(1.4)
top0 = Inches(2.0)
for i, (head, mid, body, c) in enumerate(cards):
    left = Inches(8.0)
    top = top0 + i * (card_h + Inches(0.12))
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left, top, Inches(4.8), card_h)
    box.line.color.rgb = c
    box.line.width = Pt(1.25)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    add_text(s, head, left + Inches(0.15), top + Inches(0.1),
             Inches(1.5), Inches(0.4),
             size=13, bold=True, color=c)
    add_text(s, mid, left + Inches(1.5), top + Inches(0.1),
             Inches(3.2), Inches(0.4),
             size=14, bold=True, color=NAVY)
    add_text(s, body, left + Inches(0.15), top + Inches(0.55),
             Inches(4.6), Inches(0.8),
             size=12, color=DARK)

footer(s, 7, TOTAL)

# =========================================================================
# Slide 8 — Feature Engineering
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "06  ·  FEATURES", "What We Feed the Model")

groups = [
    ("CALENDAR", TEAL,
     ["is_weekend", "dow_sin / dow_cos\n(cyclical day-of-week)"]),
    ("WEATHER", ACCENT,
     ["mean_temperature", "precipitation_sum",
      "max_wind_speed", "has_precipitation",
      "heavy_precipitation\n(top decile)", "high_wind (top decile)"]),
    ("LAG DEMAND", NAVY,
     ["lag_1_trip_count\n(yesterday)", "lag_7_trip_count\n(same weekday last week)"]),
    ("LOCATION", TEAL,
     ["start_lat / start_lng",
      "has_park / has_pier /\nhas_plaza / has_square",
      "station_cluster (one-hot)"]),
]

col_w = Inches(2.95)
gap = Inches(0.15)
top = Inches(1.95)
left0 = Inches(0.6)

for i, (head, c, items) in enumerate(groups):
    left = left0 + i * (col_w + gap)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             left, top, col_w, Inches(4.4))
    box.line.color.rgb = c
    box.line.width = Pt(1.5)
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    add_text(s, head, left, top + Inches(0.15), col_w, Inches(0.4),
             size=14, bold=True, color=c, align=PP_ALIGN.CENTER)
    body_tb = s.shapes.add_textbox(left + Inches(0.15), top + Inches(0.7),
                                   col_w - Inches(0.3), Inches(3.6))
    tf = body_tb.text_frame
    tf.word_wrap = True
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = "• " + it
        run.font.size = Pt(13)
        run.font.color.rgb = DARK
        run.font.name = "Calibri"

# Bottom callout
callout = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(0.6), Inches(6.45),
                             Inches(12.1), Inches(0.6))
callout.line.color.rgb = NAVY
callout.line.width = Pt(1.25)
callout.fill.solid(); callout.fill.fore_color.rgb = NAVY
add_text(s,
         "Target = log(1 + daily_trip_count)   ·   "
         "Same-day rider/bike shares EXCLUDED to avoid leakage",
         Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.5),
         size=15, bold=True, color=LIGHT, align=PP_ALIGN.CENTER)

footer(s, 8, TOTAL)

# =========================================================================
# Slide 9 — Models & Results
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "07  ·  MODELING", "Four Models, One Clear Winner")

# Models card
mb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), Inches(1.95),
                        Inches(5.0), Inches(4.8))
mb.line.color.rgb = TEAL; mb.line.width = Pt(1.5)
mb.fill.solid(); mb.fill.fore_color.rgb = LIGHT
add_text(s, "MODELS COMPARED", Inches(0.75), Inches(2.05),
         Inches(4.7), Inches(0.4),
         size=14, bold=True, color=TEAL)
add_bullets(s, [
    "Ridge Regression — linear baseline (α=1).",
    "Random Forest — 300 trees, depth 12, leaf 5.",
    "Gradient Boosting — 300 trees, lr 0.05, depth 3.",
    "Tuned Random Forest — 3-fold GridSearchCV\n   (n=300, depth=None, leaf=3).",
], Inches(0.75), Inches(2.55), Inches(4.7), Inches(4.0), size=15)

# Results table
header_top = Inches(2.05)
rows = [
    ("Model", "MAE", "RMSE", "R²"),
    ("Tuned Random Forest", "19.75", "28.13", "0.289"),
    ("Random Forest",       "21.23", "29.70", "0.207"),
    ("Gradient Boosting",   "24.70", "32.71", "0.038"),
    ("Ridge Regression",    "25.85", "35.78", "−0.151"),
]
tbl_left = Inches(5.95)
tbl_top = header_top
col_widths = [Inches(3.4), Inches(1.3), Inches(1.3), Inches(1.3)]
row_h = Inches(0.55)

for r_idx, row in enumerate(rows):
    x = tbl_left
    is_header = (r_idx == 0)
    is_best = (r_idx == 1)
    for c_idx, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, tbl_top + r_idx * row_h,
                                  col_widths[c_idx], row_h)
        cell.line.color.rgb = MUTED
        cell.line.width = Pt(0.5)
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = NAVY
            txt_color = LIGHT; bold = True
        elif is_best:
            cell.fill.fore_color.rgb = ACCENT
            txt_color = LIGHT; bold = True
        else:
            cell.fill.fore_color.rgb = LIGHT
            txt_color = DARK; bold = False
        align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
        add_text(s, val,
                 x + Inches(0.15), tbl_top + r_idx * row_h + Inches(0.12),
                 col_widths[c_idx] - Inches(0.2), Inches(0.4),
                 size=15, bold=bold, color=txt_color, align=align)
        x += col_widths[c_idx]

add_text(s,
         "Tuned RF wins on every metric — picked as the final model.",
         Inches(5.95), Inches(5.7), Inches(7.0), Inches(0.5),
         size=16, bold=True, color=NAVY)
add_text(s,
         "Time-based holdout — most recent 20% of dates.",
         Inches(5.95), Inches(6.15), Inches(7.0), Inches(0.5),
         size=14, color=MUTED)

footer(s, 9, TOTAL)

# =========================================================================
# Slide 10 — Why & Diagnostics
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "08  ·  WHAT THE MODEL LEARNED", "Lag Demand & Weather Drive Predictions")

# Feature importance "bars"
add_text(s, "TOP FEATURE IMPORTANCES",
         Inches(0.6), Inches(1.95), Inches(6), Inches(0.4),
         size=14, bold=True, color=TEAL)

feats = [
    ("lag_1_trip_count",   0.506),
    ("mean_temperature",   0.194),
    ("max_wind_speed",     0.075),
    ("precipitation_sum",  0.061),
    ("lag_7_trip_count",   0.056),
    ("dow_sin",            0.020),
    ("start_lat",          0.019),
    ("start_lng",          0.018),
    ("heavy_precipitation",0.017),
    ("station_cluster_0",  0.011),
]
bar_left = Inches(2.5)
bar_max_w = Inches(4.0)
bar_h = Inches(0.28)
top0 = Inches(2.5)
gap = Inches(0.10)
max_imp = max(v for _, v in feats)
for i, (name, val) in enumerate(feats):
    y = top0 + i * (bar_h + gap)
    add_text(s, name, Inches(0.6), y - Inches(0.02), Inches(1.85), bar_h,
             size=11, color=DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             bar_left, y,
                             int(bar_max_w * (val / max_imp)), bar_h)
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL if i > 0 else ACCENT
    add_text(s, f"{val:.3f}",
             bar_left + int(bar_max_w * (val / max_imp)) + Inches(0.05),
             y - Inches(0.02), Inches(0.7), bar_h,
             size=11, color=MUTED)

# Right side: error-by-quartile
add_text(s, "ERROR BY DEMAND QUARTILE",
         Inches(7.4), Inches(1.95), Inches(5.5), Inches(0.4),
         size=14, bold=True, color=TEAL)

err_rows = [
    ("Group", "Actual", "Predicted", "MAE"),
    ("Low",         "17.0", "25.2", "14.3"),
    ("Medium-Low",  "41.2", "47.0", "16.6"),
    ("Medium-High", "61.5", "60.7", "16.0"),
    ("High",        "99.9", "83.1", "32.6"),
]
e_left = Inches(7.4)
e_top = Inches(2.5)
e_widths = [Inches(1.7), Inches(1.2), Inches(1.4), Inches(1.2)]
e_h = Inches(0.45)
for r_idx, row in enumerate(err_rows):
    x = e_left
    is_header = (r_idx == 0)
    is_high = (r_idx == 4)
    for c_idx, val in enumerate(row):
        cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  x, e_top + r_idx * e_h,
                                  e_widths[c_idx], e_h)
        cell.line.color.rgb = MUTED
        cell.line.width = Pt(0.5)
        cell.fill.solid()
        if is_header:
            cell.fill.fore_color.rgb = NAVY; txt = LIGHT; bold = True
        elif is_high:
            cell.fill.fore_color.rgb = ACCENT; txt = LIGHT; bold = True
        else:
            cell.fill.fore_color.rgb = LIGHT; txt = DARK; bold = False
        align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
        add_text(s, val,
                 x + Inches(0.1), e_top + r_idx * e_h + Inches(0.08),
                 e_widths[c_idx] - Inches(0.15), Inches(0.35),
                 size=13, bold=bold, color=txt, align=align)
        x += e_widths[c_idx]

add_text(s,
         "Model is well-calibrated except for the highest-demand quartile,\n"
         "where it under-predicts (≈ −17 trips on avg). A known cost of\n"
         "log-transform + tree smoothing on heavy-tailed targets.",
         Inches(7.4), Inches(5.0), Inches(5.5), Inches(2.0),
         size=13, color=DARK)

footer(s, 10, TOTAL)

# =========================================================================
# Slide 11 — Limits & Future Work
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s)
section_header(s, "09  ·  HONEST", "Limitations & What We'd Do Next")

# Two columns: limitations / future work
def column(title, color, items, left):
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               left, Inches(1.95),
                               Inches(6.0), Inches(4.8))
    panel.line.color.rgb = color
    panel.line.width = Pt(1.5)
    panel.fill.solid(); panel.fill.fore_color.rgb = LIGHT
    add_text(s, title, left + Inches(0.25), Inches(2.1),
             Inches(5.5), Inches(0.5),
             size=18, bold=True, color=color)
    add_bullets(s, items, left + Inches(0.25), Inches(2.7),
                Inches(5.5), Inches(4.0), size=15)

column("LIMITATIONS", ACCENT, [
    "Only 2 months of data — no seasonal cycle.",
    "Under-prediction at peak-demand stations.",
    "No holiday or special-event indicators.",
    "Weather is daily — sub-day patterns lost.",
], Inches(0.6))

column("NEXT STEPS", TEAL, [
    "Span a full year to capture seasonality.",
    "Quantile regression / two-stage surge model\n   for peak days.",
    "Add holidays, NYC events, transit disruptions.",
    "Deploy as a Streamlit station-day demand dashboard.",
], Inches(6.85))

footer(s, 11, TOTAL)

# =========================================================================
# Slide 12 — Closing & Contributions
# =========================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, NAVY)

bar1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                          0, Inches(6.6), SW, Inches(0.12))
bar1.line.fill.background()
bar1.fill.solid(); bar1.fill.fore_color.rgb = ACCENT

add_text(s, "10  ·  WRAP-UP",
         Inches(0.7), Inches(0.6), Inches(12), Inches(0.4),
         size=14, bold=True, color=ACCENT)

add_text(s, "Thank you  ·  Questions?",
         Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
         size=42, bold=True, color=LIGHT)

# Headline takeaway
take = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.7), Inches(2.3),
                          Inches(11.9), Inches(1.5))
take.line.color.rgb = ACCENT
take.line.width = Pt(1.5)
take.fill.solid(); take.fill.fore_color.rgb = NAVY
add_text(s, "KEY TAKEAWAY",
         Inches(1.0), Inches(2.45), Inches(11.5), Inches(0.4),
         size=13, bold=True, color=ACCENT)
add_text(s,
         "A tuned Random Forest on lag demand + weather + clusters\n"
         "predicts station-day Citi Bike demand with MAE ≈ 20 trips\n"
         "on a true future holdout — strong enough to support rebalancing decisions.",
         Inches(1.0), Inches(2.9), Inches(11.5), Inches(1.0),
         size=15, color=LIGHT)

# Contributions
add_text(s, "TEAM CONTRIBUTIONS",
         Inches(0.7), Inches(4.1), Inches(12), Inches(0.4),
         size=14, bold=True, color=ACCENT)

contribs = [
    ("Helena Li (yl6029)",     "Data acquisition & cleaning · station-day analytic table"),
    ("Wentao Zhong (wz2753)",  "Exploratory data analysis · figures & summary tables"),
    ("Kevin Ma (km4189)",      "Clustering · feature engineering · supervised models · evaluation"),
    ("Ketaki Dabade (kvd2112)","Final report · README · evaluation interpretation · slides"),
]
for i, (name, role) in enumerate(contribs):
    y = Inches(4.55) + i * Inches(0.42)
    add_text(s, name, Inches(0.9), y, Inches(3.6), Inches(0.4),
             size=14, bold=True, color=LIGHT)
    add_text(s, role, Inches(4.6), y, Inches(8.0), Inches(0.4),
             size=14, color=LIGHT)

add_text(s, "github.com/Helenaliyz/5243Project4",
         Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         size=12, color=ACCENT)

footer(s, 12, TOTAL)

# ---- Save ----
prs.save(str(OUT))
print(f"Wrote {OUT}")
